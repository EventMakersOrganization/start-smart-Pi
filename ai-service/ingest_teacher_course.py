"""
Ingest teacher's "Programmation Procedurale 1" course folder into MongoDB + ChromaDB.

Parses the fiche module for chapter metadata, extracts course content from
PPTX/PDF files, parses quiz DOCX files into MCQ exercises, then:
  1. Removes **only** prior "Programmation Procédurale 1" rows (Mongo + matching
     Chroma chunks), leaving other subjects untouched
  2. Inserts 8 courses (one per chapter) with version-based subchapters
  3. Inserts all parsed MCQ exercises linked to their parent course
  4. Triggers the embedding pipeline for RAG

Usage:
    cd ai-service
    python ingest_teacher_course.py

Dangerous legacy options (full DB wipe — use only if you intend to erase everything):

    python ingest_teacher_course.py --wipe-all-mongo --wipe-all-chroma

Environment:
    TEACHER_COURSE_ROOT — root folder (e.g. ``…\\Support_Cours_Préparation``). Each
    ``chapter/X.Y/Cours/*.pptx|pdf`` is one Mongo module; files are copied to
    ``backend/uploads/subjects/cours/`` and ``fileUrl`` is set so the student UI
    shows separate files like other subjects. Text is still extracted for embeddings.

    INGEST_DEFAULT_INSTRUCTOR_ID — optional MongoDB ObjectId (hex) of the User who
    should own inserted courses. If unset, the first user with role instructor
    (or teacher) in ``users`` is used. If none is found, courses are inserted
    without ``instructorId`` (same as before).
"""
from __future__ import annotations

import argparse
import glob
import io
import logging
import os
import re
import shutil
import sys
import time
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

_SERVICE_ROOT = Path(__file__).resolve().parent
if str(_SERVICE_ROOT) not in sys.path:
    sys.path.insert(0, str(_SERVICE_ROOT))

from bson import ObjectId
from core import config

try:
    from pymongo import MongoClient
except ImportError:
    raise ImportError("pymongo is required")

logging.basicConfig(level=logging.INFO, format="%(levelname)s | %(message)s")
log = logging.getLogger("ingest")

_STARTSMART_ROOT = Path(__file__).resolve().parent.parent
_UPLOADS_SUBJECTS_COURS = _STARTSMART_ROOT / "backend" / "uploads" / "subjects" / "cours"

TEACHER_ROOT = Path(
    os.getenv(
        "TEACHER_COURSE_ROOT",
        r"C:\Users\kmarb\Downloads\pi\Support_Cours_Préparation",
    )
)


def _resolve_course_instructor_id(db: Any) -> ObjectId | None:
    """Match Nest-created courses: set ``instructorId`` when possible."""
    raw = os.getenv("INGEST_DEFAULT_INSTRUCTOR_ID", "").strip()
    if raw:
        if ObjectId.is_valid(raw):
            oid = ObjectId(raw)
            if db.users.find_one({"_id": oid}):
                return oid
            log.warning(
                "INGEST_DEFAULT_INSTRUCTOR_ID=%s not found in users; trying first instructor",
                raw,
            )
        else:
            log.warning(
                "INGEST_DEFAULT_INSTRUCTOR_ID is not a valid ObjectId; trying first instructor",
            )
    user = db.users.find_one(
        {"role": {"$regex": r"^(instructor|teacher)$", "$options": "i"}}
    )
    if user:
        return user["_id"]
    log.warning(
        "No instructor user in DB; courses will have no instructorId "
        "(set INGEST_DEFAULT_INSTRUCTOR_ID to a valid User _id)",
    )
    return None


def _course_document_base(
    *,
    instructor_id: ObjectId | None,
    now: datetime,
) -> dict[str, Any]:
    """Fields aligned with Mongoose ``Course`` (timestamps + version key)."""
    doc: dict[str, Any] = {
        "createdAt": now,
        "updatedAt": now,
        "__v": 0,
    }
    if instructor_id is not None:
        doc["instructorId"] = instructor_id
    return doc


def _programmation_subject_mongo_filter() -> dict[str, Any]:
    """Matches ``subject`` on course/exercise docs for this ingest."""
    return {
        "subject": {"$regex": r"Programmation\s+Proc[ée]durale\s*1", "$options": "i"},
    }


def _delete_chroma_chunks_for_course_ids(course_ids: list[str]) -> None:
    """Remove only vectors tied to these Mongo course ids (other subjects stay)."""
    if not course_ids:
        return
    try:
        from core.chroma_setup import get_or_create_collection
    except ImportError:
        log.warning("ChromaDB not available; skipping chunk delete.")
        return
    for col_name in ("course_embeddings", "course_chunks"):
        try:
            coll = get_or_create_collection(col_name)
            try:
                coll.delete(where={"course_id": {"$in": course_ids}})
            except Exception:
                for cid in course_ids:
                    try:
                        coll.delete(where={"course_id": cid})
                    except Exception:
                        pass
            log.info(
                "ChromaDB %s: removed chunks for %d course id(s).",
                col_name,
                len(course_ids),
            )
        except Exception as e:
            log.warning("ChromaDB %s prune skipped: %s", col_name, e)


def _remove_prior_programmation_data(
    db: Any,
    courses_coll: Any,
    exercises_coll: Any,
    *,
    wipe_all_mongo: bool,
    wipe_all_chroma: bool,
) -> None:
    """
    Default: delete only Programmation Procédurale 1 courses/exercises and matching
    Chroma chunks. Optional flags restore the old full-database wipe.
    """
    try:
        from core.chroma_setup import delete_collection
    except ImportError:
        delete_collection = None  # type: ignore

    if wipe_all_mongo:
        del_c = courses_coll.delete_many({})
        del_e = exercises_coll.delete_many({})
        log.warning(
            "Removed ALL MongoDB courses (%d) and exercises (%d).",
            del_c.deleted_count,
            del_e.deleted_count,
        )
        if wipe_all_chroma and delete_collection:
            for col_name in ("course_embeddings", "course_chunks"):
                try:
                    delete_collection(col_name)
                    log.warning("Cleared entire ChromaDB collection: %s", col_name)
                except Exception as e:
                    log.warning("ChromaDB full delete %s: %s", col_name, e)
        elif wipe_all_chroma:
            log.warning("--wipe-all-chroma ignored (chroma_setup unavailable).")
        return

    filt = _programmation_subject_mongo_filter()
    old_docs = list(courses_coll.find(filt, {"_id": 1}))
    old_ids = [str(d["_id"]) for d in old_docs]

    if old_ids:
        if wipe_all_chroma and delete_collection:
            for col_name in ("course_embeddings", "course_chunks"):
                try:
                    delete_collection(col_name)
                    log.warning("Cleared entire ChromaDB collection: %s", col_name)
                except Exception as e:
                    log.warning("ChromaDB full delete %s: %s", col_name, e)
        else:
            _delete_chroma_chunks_for_course_ids(old_ids)

        obj_ids: list[ObjectId] = []
        for x in old_ids:
            if ObjectId.is_valid(x):
                obj_ids.append(ObjectId(x))

        if obj_ids:
            ex_res = exercises_coll.delete_many(
                {"$or": [filt, {"courseId": {"$in": obj_ids}}]}
            )
        else:
            ex_res = exercises_coll.delete_many(filt)
        c_res = courses_coll.delete_many(filt)
        log.info(
            "Replaced prior Programmation data: %d course(s), %d exercise(s) (Mongo).",
            c_res.deleted_count,
            ex_res.deleted_count,
        )
    else:
        log.info("No existing Programmation courses in MongoDB to replace.")
        if wipe_all_chroma and delete_collection:
            for col_name in ("course_embeddings", "course_chunks"):
                try:
                    delete_collection(col_name)
                    log.warning("Cleared entire ChromaDB collection: %s", col_name)
                except Exception as e:
                    log.warning("ChromaDB full delete %s: %s", col_name, e)


# ── Chapter metadata extracted from fiche module ─────────────────────────
# folder_num -> (title, description/objectives)
CHAPTERS: dict[int, dict[str, str]] = {
    0: {
        "title": "Chapitre 0 : Chapitre introductif",
        "description": "Mémoriser des définitions relatives à l'algorithmique et à la programmation C.",
    },
    1: {
        "title": "Chapitre 1 : Lecture & Écriture des données",
        "description": (
            "Lister les étapes pour passer d'un code source à un exécutable. "
            "Décrire la structure d'un programme C. Utiliser les variables. "
            "Examiner la portée d'une variable. Pratiquer les fonctions d'entrées/sorties. "
            "Utiliser et ordonner les opérateurs arithmétiques et logiques."
        ),
    },
    2: {
        "title": "Chapitre 2 : Les structures conditionnelles",
        "description": (
            "Appliquer la structure if else. Appliquer la structure switch. "
            "Distinguer les structures conditionnelles (if, if else et switch) du langage C."
        ),
    },
    3: {
        "title": "Chapitre 3 : Les structures itératives",
        "description": (
            "Utiliser la boucle for. Utiliser la boucle do while. "
            "Utiliser la boucle while. Différencier les structures itératives."
        ),
    },
    4: {
        "title": "Chapitre 4 : Les tableaux unidimensionnels et bidimensionnels",
        "description": (
            "Définir un tableau. Pratiquer l'ajout, la suppression et le parcours. "
            "Appliquer la recherche séquentielle et dichotomique. Ordonner un tableau. "
            "Identifier les types de tableaux. Parcourir un tableau bidimensionnel."
        ),
    },
    5: {
        "title": "Chapitre 5 : Les chaînes de caractères",
        "description": (
            "Définir une chaîne de caractères. Pratiquer les fonctions de saisie et d'affichage. "
            "Utiliser les fonctions de la bibliothèque string.h. "
            "Pratiquer les tableaux de chaînes de caractères."
        ),
    },
    6: {
        "title": "Chapitre 6 : Les structures",
        "description": (
            "Définir une structure avec des champs simples, des champs de type tableau "
            "et des champs de type enregistrement. Utiliser une variable de type structure. "
            "Utiliser un tableau de structures."
        ),
    },
    7: {
        "title": "Chapitre 7 : Les pointeurs",
        "description": (
            "Définir un pointeur. Utiliser les pointeurs dans les prototypes des fonctions "
            "et dans les sous-programmes. Appliquer les pointeurs pour la manipulation des tableaux."
        ),
    },
    8: {
        "title": "Chapitre 8 : Les fonctions",
        "description": (
            "Définir une fonction, ses intérêts et ses caractéristiques. "
            "Représenter le prototype d'une fonction. Pratiquer l'appel des fonctions. "
            "Différencier les paramètres effectifs des paramètres formels. "
            "Distinguer les modes de passage des paramètres. "
            "Utiliser les tableaux dans des fonctions."
        ),
    },
}

# Version descriptions (X.Y -> short topic) from the fiche module
VERSION_TITLES: dict[str, str] = {
    "1.1": "Étapes pour passer d'un code source à un exécutable",
    "1.2": "Structure d'un programme C",
    "1.3": "Les variables",
    "1.4": "Portée d'une variable",
    "1.5": "Fonctions d'entrées / sorties",
    "1.6": "Opérateurs arithmétiques et logiques",
    "1.7": "Priorité des opérateurs",
    "2.1": "Structure if else",
    "2.2": "Structure switch",
    "2.3": "Distinction des structures conditionnelles",
    "3.1": "Boucle for",
    "3.2": "Boucle do while",
    "3.3": "Boucle while",
    "3.4": "Différenciation des structures itératives",
    "4.1": "Définition d'un tableau",
    "4.2": "Opération d'ajout",
    "4.3": "Opération de suppression",
    "4.4": "Opération de parcours",
    "4.5": "Recherche séquentielle",
    "4.6": "Tri d'un tableau",
    "4.7": "Recherche dichotomique",
    "4.8": "Types des tableaux",
    "4.9": "Parcours d'un tableau bidimensionnel",
    "5.1": "Définition d'une chaîne de caractères",
    "5.2": "Saisie et affichage d'une chaîne",
    "5.3": "Fonctions de string.h",
    "5.4": "Tableaux de chaînes de caractères",
    "5.5": "Manipulation avancée de chaînes",
    "5.6": "Opérations sur chaînes",
    "5.7": "Exercices pratiques sur les chaînes",
    "6.1": "Définition d'une structure",
    "6.2": "Variable de type structure",
    "6.3": "Tableau de structures",
    "7.1": "Définition d'un pointeur",
    "7.2": "Pointeurs et fonctions",
    "7.3": "Pointeurs et tableaux",
    "8.1": "Définition d'une fonction",
    "8.2": "Prototype d'une fonction",
    "8.3": "Appel de fonctions",
    "8.4": "Paramètres effectifs vs formels",
    "8.5": "Modes de passage des paramètres",
    "8.6": "Tableaux dans les fonctions",
}


# ── File content extractors ──────────────────────────────────────────────

def extract_pptx_text(path: Path) -> str:
    """Extract all text from a PPTX presentation."""
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for slide in list(prs.slides):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_texts.append(t)
        if slide_texts:
            parts.append("\n".join(slide_texts))
    return "\n\n".join(parts)


def extract_pdf_text(path: Path) -> str:
    """Extract text from a PDF file."""
    try:
        from pypdf import PdfReader
    except ImportError:
        import PyPDF2
        reader = PyPDF2.PdfReader(str(path))
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    reader = PdfReader(str(path))
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def extract_course_text(cours_dir: Path) -> str:
    """Extract text from the best available file in a Cours/ folder.
    Prefers PPTX (richer text), falls back to PDF."""
    pptx_files = list(cours_dir.glob("*.pptx"))
    pdf_files = list(cours_dir.glob("*.pdf"))

    texts: list[str] = []
    if pptx_files:
        for f in sorted(pptx_files):
            try:
                texts.append(extract_pptx_text(f))
            except Exception as e:
                log.warning("  PPTX extraction failed for %s: %s", f.name, e)
    elif pdf_files:
        for f in sorted(pdf_files):
            try:
                texts.append(extract_pdf_text(f))
            except Exception as e:
                log.warning("  PDF extraction failed for %s: %s", f.name, e)

    return "\n\n".join(t for t in texts if t.strip())


def _copy_module_file_to_uploads(
    src: Path, chapter_num: int, version: str
) -> tuple[str, str]:
    """
    Copy a PPTX/PDF from the teacher folder into backend/uploads/subjects/cours
    so the student app can open it like other subjects. Returns (relativeUrl, storedFileName).
    """
    if not src.is_file():
        return "", ""
    _UPLOADS_SUBJECTS_COURS.mkdir(parents=True, exist_ok=True)
    safe_stem = re.sub(r"[^\w\-\s]+", "_", src.stem, flags=re.UNICODE)[:80].strip().replace(" ", "_")
    ext = (src.suffix or "").lower() or ".bin"
    base_name = f"pp1_ch{chapter_num}_v{version.replace('.', '_')}_{safe_stem}{ext}"
    dest = _UPLOADS_SUBJECTS_COURS / base_name
    if dest.exists():
        dest = _UPLOADS_SUBJECTS_COURS / f"{dest.stem}_{int(time.time())}{ext}"
    shutil.copy2(src, dest)
    return f"/uploads/subjects/cours/{dest.name}", dest.name


def _list_module_files_with_text(cours_dir: Path) -> list[dict[str, Any]]:
    """One entry per PPTX/PDF in Cours (same structure as on disk)."""
    out: list[dict[str, Any]] = []
    files = sorted(
        list(cours_dir.glob("*.pptx")) + list(cours_dir.glob("*.pdf")),
        key=lambda p: p.name.lower(),
    )
    for f in files:
        try:
            if f.suffix.lower() == ".pptx":
                t = extract_pptx_text(f)
            else:
                t = extract_pdf_text(f)
            out.append({"path": f, "text": t})
        except Exception as e:
            log.warning("  File extract failed %s: %s", f.name, e)
    return out


# ── Quiz DOCX parser ─────────────────────────────────────────────────────

_OPTION_LINE = re.compile(r"^\s*([A-Da-d])\s*[.)]\s*(.+)", re.DOTALL)
_INLINE_OPTIONS = re.compile(
    r"[Aa]\s*[.)]\s*(.+?)\s*[Bb]\s*[.)]\s*(.+?)\s*[Cc]\s*[.)]\s*(.+?)\s*[Dd]\s*[.)]\s*(.+)"
)
_CORRECT_MARKER = re.compile(
    r"(✅|bonne\s+r[ée]ponse|r[ée]ponse\s*:)\s*[:]?\s*([A-Da-d])\s*[.)]\s*(.+)",
    re.IGNORECASE,
)


def parse_quiz_docx(path: Path) -> list[dict[str, Any]]:
    """Parse a quiz DOCX into a list of MCQ dicts."""
    import docx

    doc = docx.Document(str(path))
    raw_lines = [p.text for p in doc.paragraphs]

    questions: list[dict[str, Any]] = []
    buf_question: list[str] = []
    buf_options: list[str] = []
    correct: str | None = None

    def _flush():
        nonlocal buf_question, buf_options, correct
        q_text = "\n".join(buf_question).strip()
        if not q_text or not buf_options:
            buf_question, buf_options, correct = [], [], None
            return
        q_text = re.sub(r"^(Question\s*:?\s*)", "", q_text, flags=re.IGNORECASE).strip()
        if not correct and buf_options:
            correct = buf_options[0]
        questions.append({
            "type": "MCQ",
            "content": q_text,
            "options": list(buf_options),
            "correctAnswer": correct or buf_options[0],
            "difficulty": "medium",
        })
        buf_question, buf_options, correct = [], [], None

    for line in raw_lines:
        line_stripped = line.strip()
        if not line_stripped:
            if buf_question and buf_options:
                _flush()
            continue

        cm = _CORRECT_MARKER.search(line_stripped)
        if cm:
            correct = f"{cm.group(2).upper()}. {cm.group(3).strip()}"
            if buf_question and buf_options:
                _flush()
            continue

        im = _INLINE_OPTIONS.search(line_stripped)
        if im:
            before_a = line_stripped[:im.start()].strip()
            if before_a:
                q_part = re.sub(r"^(Question\s*:?\s*)", "", before_a, flags=re.IGNORECASE).strip()
                if q_part:
                    buf_question.append(q_part)
            buf_options = [
                f"A. {im.group(1).strip()}",
                f"B. {im.group(2).strip()}",
                f"C. {im.group(3).strip()}",
                f"D. {im.group(4).strip()}",
            ]
            continue

        om = _OPTION_LINE.match(line_stripped)
        if om:
            buf_options.append(f"{om.group(1).upper()}. {om.group(2).strip()}")
            continue

        buf_question.append(line_stripped)

    if buf_question and buf_options:
        _flush()

    return questions


# ── Walk the teacher's folder ─────────────────────────────────────────────

def discover_versions(chapter_dir: Path) -> list[dict[str, Any]]:
    """Discover all X.Y version subfolders under a chapter directory."""
    versions: list[dict[str, Any]] = []
    for child in sorted(chapter_dir.iterdir()):
        if not child.is_dir():
            continue
        if not re.match(r"\d+\.\d+$", child.name):
            continue
        v: dict[str, Any] = {"version": child.name, "path": child, "text": "", "quizzes": []}

        cours_dir = None
        for candidate in ("Cours", "cours"):
            d = child / candidate
            if d.is_dir():
                cours_dir = d
                break
        if cours_dir:
            v["module_files"] = _list_module_files_with_text(cours_dir)
            v["text"] = "\n\n".join(
                (x.get("text") or "").strip() for x in v["module_files"] if x.get("text")
            )
            if not v["text"]:
                v["text"] = extract_course_text(cours_dir)

        quiz_dir = None
        for candidate in ("Quizz", "Quiz", "quizz", "quiz"):
            d = child / candidate
            if d.is_dir():
                quiz_dir = d
                break
        if quiz_dir:
            for qf in sorted(quiz_dir.glob("*.docx")):
                try:
                    v["quizzes"].extend(parse_quiz_docx(qf))
                except Exception as e:
                    log.warning("  Quiz parse error %s: %s", qf.name, e)

        versions.append(v)
    return versions


def discover_extra_quizzes(chapter_dir: Path) -> list[dict[str, Any]]:
    """Find quizzes in non-version subfolders (e.g. 'Appliquer Structures iteratives')."""
    extras: list[dict[str, Any]] = []
    for child in chapter_dir.iterdir():
        if not child.is_dir():
            continue
        if re.match(r"\d+\.\d+$", child.name):
            continue
        for qf in sorted(child.glob("*.docx")):
            if "quiz" in qf.stem.lower():
                try:
                    extras.extend(parse_quiz_docx(qf))
                except Exception as e:
                    log.warning("  Extra quiz parse error %s: %s", qf.name, e)
    return extras


# ── Main ingestion ────────────────────────────────────────────────────────

def run_ingest(
    *,
    wipe_all_mongo: bool = False,
    wipe_all_chroma: bool = False,
) -> dict[str, Any]:
    log.info("=" * 60)
    log.info("  Teacher Course Ingestion")
    log.info("=" * 60)
    log.info("Source: %s", TEACHER_ROOT)

    if not TEACHER_ROOT.is_dir():
        log.error("Teacher folder not found: %s", TEACHER_ROOT)
        sys.exit(1)

    # ── Connect to MongoDB ────────────────────────────────────────────
    client = MongoClient(config.MONGODB_URI)
    db = client[config.MONGODB_DB_NAME]
    courses_coll = db["courses"]
    exercises_coll = db["exercises"]
    log.info("Connected to MongoDB: %s", config.MONGODB_DB_NAME)

    instructor_oid = _resolve_course_instructor_id(db)
    if instructor_oid is not None:
        log.info("Resolved course instructorId: %s", instructor_oid)

    _remove_prior_programmation_data(
        db,
        courses_coll,
        exercises_coll,
        wipe_all_mongo=wipe_all_mongo,
        wipe_all_chroma=wipe_all_chroma,
    )

    # ── Process chapters 1-8 ──────────────────────────────────────────
    course_ids: list[str] = []
    total_exercises = 0

    for ch_num in range(1, 9):
        ch_dir = TEACHER_ROOT / str(ch_num)
        meta = CHAPTERS.get(ch_num, {
            "title": f"Chapitre {ch_num}",
            "description": "",
        })

        log.info("")
        log.info("─── %s ───", meta["title"])

        if not ch_dir.is_dir():
            log.warning("  Chapter folder %d not found, skipping", ch_num)
            continue

        versions = discover_versions(ch_dir)
        log.info("  Found %d versions", len(versions))

        sub_chapters: list[dict[str, Any]] = []
        all_quizzes: list[dict[str, Any]] = []
        mod_order = 0

        for v in versions:
            ver = v["version"]
            topic = VERSION_TITLES.get(ver, f"Section {ver}")
            base_title = f"{ver} - {topic}"
            mfiles = v.get("module_files") or []
            all_quizzes.extend(v["quizzes"])

            if mfiles:
                for mf in mfiles:
                    p = mf.get("path")
                    txt = (mf.get("text") or "").strip()
                    rel_url, stored = ("", "")
                    if isinstance(p, Path) and p.is_file():
                        rel_url, stored = _copy_module_file_to_uploads(p, ch_num, ver)
                        log.info(
                            "  [%s] %s (%d chars) -> %s",
                            ver,
                            p.name,
                            len(txt),
                            rel_url or "(copy failed)",
                        )
                    else:
                        log.info("  [%s] skip file (invalid path)", ver)

                    title = (
                        base_title
                        if len(mfiles) == 1
                        else f"{base_title} — {getattr(p, 'stem', 'fichier')}"
                    )
                    row: dict[str, Any] = {
                        "title": title,
                        "description": txt,
                        "order": mod_order,
                        "contents": [],
                    }
                    if rel_url:
                        row["fileUrl"] = rel_url
                        row["fileName"] = stored or (
                            p.name if isinstance(p, Path) else "cours"
                        )
                    sub_chapters.append(row)
                    mod_order += 1
            else:
                text = (v.get("text") or "").strip()
                if text:
                    log.info("  [%s] %d chars (combined)", ver, len(text))
                else:
                    log.info("  [%s] no course content", ver)
                sub_chapters.append(
                    {
                        "title": base_title,
                        "description": text,
                        "order": mod_order,
                        "contents": [],
                    }
                )
                mod_order += 1

        extra_q = discover_extra_quizzes(ch_dir)
        all_quizzes.extend(extra_q)
        if extra_q:
            log.info("  +%d extra quizzes from subfolders", len(extra_q))

        now = datetime.now(timezone.utc)
        course_doc = {
            "title": meta["title"],
            "description": meta["description"],
            "subChapters": sub_chapters,
            "level": "Beginner",
            "subject": "Programmation Procédurale 1",
            **_course_document_base(instructor_id=instructor_oid, now=now),
        }
        result = courses_coll.insert_one(course_doc)
        cid = str(result.inserted_id)
        course_ids.append(cid)
        log.info("  => Course inserted: %s (id=%s)", meta["title"], cid)

        for q in all_quizzes:
            q["courseId"] = ObjectId(cid)
            q["subject"] = "Programmation Procédurale 1"
            q["topic"] = meta["title"]
        if all_quizzes:
            exercises_coll.insert_many(all_quizzes)
            total_exercises += len(all_quizzes)
            log.info("  => %d exercises inserted", len(all_quizzes))
        else:
            log.info("  => 0 exercises (no quizzes found)")

    # ── Trigger embedding pipeline ────────────────────────────────────
    log.info("")
    log.info("=" * 60)
    log.info("  Embedding courses into ChromaDB")
    log.info("=" * 60)

    try:
        from embeddings.embeddings_pipeline_v2 import process_and_embed_course_chunks

        embedded = 0
        for cid in course_ids:
            ok, n_chunks = process_and_embed_course_chunks(cid)
            if ok:
                log.info("  Embedded course %s -> %d chunks", cid, n_chunks)
                embedded += 1
            else:
                log.warning("  Embedding failed for course %s", cid)
        log.info("Embedded %d / %d courses", embedded, len(course_ids))
    except Exception as e:
        log.warning("Embedding step skipped (Ollama may not be running): %s", e)
        log.info("You can embed later with: POST /embed-courses or POST /batch-embed-courses")

    # ── Summary ───────────────────────────────────────────────────────
    log.info("")
    log.info("=" * 60)
    log.info("  DONE")
    log.info("=" * 60)
    log.info("Courses inserted: %d", len(course_ids))
    for i, cid in enumerate(course_ids):
        ch = CHAPTERS.get(i + 1, {})
        log.info("  %d. %s  ->  %s", i + 1, ch.get("title", "?"), cid)
    log.info("Exercises inserted: %d", total_exercises)
    log.info("Course IDs: %s", course_ids)

    return {"course_ids": course_ids, "total_exercises": total_exercises}


if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="Ingest Programmation Procédurale 1 teacher folder into MongoDB + ChromaDB.",
    )
    parser.add_argument(
        "--wipe-all-mongo",
        action="store_true",
        help="Delete ALL documents in courses and exercises collections (dangerous).",
    )
    parser.add_argument(
        "--wipe-all-chroma",
        action="store_true",
        help="Delete entire course_embeddings and course_chunks Chroma collections.",
    )
    ns = parser.parse_args()
    run_ingest(wipe_all_mongo=ns.wipe_all_mongo, wipe_all_chroma=ns.wipe_all_chroma)
