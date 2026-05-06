"""
Intelligent document chunking for RAG.
Uses LangChain text splitters and custom logic for courses/exercises.
"""
import logging
import os
import re
from typing import Any
from urllib.parse import urlparse

try:
    from langchain_text_splitters import RecursiveCharacterTextSplitter, CharacterTextSplitter
except ImportError:
    try:
        from langchain.text_splitter import RecursiveCharacterTextSplitter, CharacterTextSplitter
    except ImportError:
        RecursiveCharacterTextSplitter = None
        CharacterTextSplitter = None

logger = logging.getLogger("document_chunker")
logger.setLevel(logging.INFO)
if not logger.handlers:
    _h = logging.StreamHandler()
    _h.setFormatter(logging.Formatter("%(asctime)s - %(levelname)s - %(message)s"))
    logger.addHandler(_h)


def _normalize_extracted_text(text: str) -> str:
    if not isinstance(text, str):
        return ""
    # Common PDF ligatures and normalization artifacts.
    text = (
        text.replace("\ufb00", "ff")
        .replace("\ufb01", "fi")
        .replace("\ufb02", "fl")
        .replace("\ufb03", "ffi")
        .replace("\ufb04", "ffl")
        .replace("\u00a0", " ")
    )
    # Many scanned/exported PDFs use U+00B4 (acute accent) separated from letters.
    text = re.sub(r"\u00B4\s*e", "é", text)
    text = re.sub(r"\u00B4\s*E", "É", text)
    text = re.sub(r"\u00B4\s*a", "á", text)
    text = re.sub(r"\u00B4\s*i", "í", text)
    text = re.sub(r"\u00B4\s*o", "ó", text)
    text = re.sub(r"\u00B4\s*u", "ú", text)
    # Heuristic mojibake repair pass for mixed UTF-8/Latin artifacts.
    # Some PDFs return strings where accented glyphs are split as "x� y".
    text = re.sub(r"([A-Za-z])\s*�\s*([A-Za-z])", r"\1\2", text)
    # Normalize common French grave/accent placeholders.
    text = text.replace("` e", "è").replace("` a", "à").replace("` u", "ù")
    text = text.replace("` E", "È").replace("` A", "À").replace("` U", "Ù")
    # Target frequent OCR-ish noise around French words.
    fixes = {
        "securite": "sécurité",
        "donnees": "données",
        "privileges": "privilèges",
        "proteger": "protéger",
        "revoquer": "révoquer",
        "necessaires": "nécessaires",
        "administration": "administration",
        "authentification": "authentification",
        "execution": "exécution",
        "operations": "opérations",
        "acceder": "accéder",
        "creer": "créer",
        "etre": "être",
    }
    lowered = text.lower()
    for raw, fixed in fixes.items():
        # Replace whole words only, case-insensitive.
        pattern = re.compile(rf"\b{re.escape(raw)}\b", re.IGNORECASE)
        text = pattern.sub(lambda m: fixed if m.group(0).islower() else fixed.capitalize(), text)
    # Remove unrecoverable replacement chars.
    text = text.replace("�", "")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _resolve_local_upload_path(file_url: str) -> str:
    """
    Convert URL like http://localhost:3000/uploads/subjects/cours/x.pdf
    to local backend path.
    """
    if not file_url:
        return ""
    parsed = urlparse(str(file_url))
    path = parsed.path if parsed.scheme else str(file_url)
    path = path.replace("\\", "/")
    marker = "/uploads/"
    if marker not in path:
        return ""
    rel = path.split(marker, 1)[1]
    service_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))  # ai-service/
    workspace_root = os.path.dirname(service_root)  # start-smart-Pi/
    full = os.path.join(workspace_root, "backend", "uploads", rel)
    return full


def _extract_text_from_local_file(path: str) -> str:
    p = str(path or "").strip()
    if not p or not os.path.exists(p):
        return ""
    ext = os.path.splitext(p)[1].lower()
    try:
        if ext == ".pdf":
            try:
                from pypdf import PdfReader  # type: ignore
            except Exception:
                from PyPDF2 import PdfReader  # type: ignore
            reader = PdfReader(p)
            raw = "\n".join((pg.extract_text() or "") for pg in reader.pages)
            return _normalize_extracted_text(raw)
        if ext in {".ppt", ".pptx"}:
            try:
                from pptx import Presentation  # type: ignore
            except Exception:
                return ""
            prs = Presentation(p)
            parts: list[str] = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    txt = getattr(shape, "text", "")
                    if txt:
                        parts.append(str(txt))
            return _normalize_extracted_text("\n".join(parts))
        if ext == ".docx":
            try:
                from docx import Document  # type: ignore
            except Exception:
                return ""
            doc = Document(p)
            raw = "\n".join(par.text for par in doc.paragraphs if par.text)
            return _normalize_extracted_text(raw)
        if ext in {".txt", ".md"}:
            with open(p, "r", encoding="utf-8", errors="ignore") as f:
                return _normalize_extracted_text(f.read())
    except Exception as e:
        logger.warning("file text extraction failed (%s): %s", p, e)
        return ""
    return ""


def _extract_subchapter_real_text(subchapter: dict[str, Any]) -> str:
    """Extract and combine text from linked content files under a subchapter."""
    if not isinstance(subchapter, dict):
        return ""
    contents = subchapter.get("contents") or []
    if not isinstance(contents, list):
        return ""
    texts: list[str] = []
    for item in contents:
        if not isinstance(item, dict):
            continue
        u = str(item.get("url") or "").strip()
        if not u:
            continue
        local_path = _resolve_local_upload_path(u)
        extracted = _extract_text_from_local_file(local_path)
        if extracted:
            texts.append(extracted[:12000])
    return _normalize_extracted_text("\n\n".join(texts))


def chunk_text_recursive(
    text: str,
    chunk_size: int = 500,
    chunk_overlap: int = 50,
) -> list[str]:
    """
    Splits text using RecursiveCharacterTextSplitter (paragraphs -> sentences -> chars).
    Best for general documents.
    """
    if not text or not str(text).strip():
        logger.warning("chunk_text_recursive: empty text")
        return []
    try:
        if RecursiveCharacterTextSplitter is None:
            logger.warning("RecursiveCharacterTextSplitter not available; using simple split")
            return _chunk_fallback(text, chunk_size, chunk_overlap)
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=len,
            separators=["\n\n", "\n", ". ", " ", ""],
        )
        chunks = splitter.split_text(text.strip())
        logger.info("chunk_text_recursive: produced %s chunks", len(chunks))
        return chunks
    except Exception as e:
        logger.error("chunk_text_recursive error: %s", e)
        return _chunk_fallback(text, chunk_size, chunk_overlap)


def _chunk_fallback(text: str, chunk_size: int, chunk_overlap: int) -> list[str]:
    """Simple overlap chunking when LangChain is not available."""
    text = text.strip()
    if not text:
        return []
    chunks = []
    start = 0
    while start < len(text):
        end = min(start + chunk_size, len(text))
        chunks.append(text[start:end])
        start = end - chunk_overlap if (end < len(text)) else len(text)
    return chunks


def chunk_by_sentences(text: str, max_sentences: int = 5) -> list[str]:
    """
    Splits text into chunks of N sentences. Preserves sentence boundaries.
    Good for course content that should stay contextual.
    """
    if not text or not str(text).strip():
        logger.warning("chunk_by_sentences: empty text")
        return []
    # Split on sentence boundaries (. ! ? followed by space or end)
    raw = re.split(r'(?<=[.!?])\s+', text.strip())
    sentences = [s.strip() for s in raw if s.strip()]
    if not sentences:
        return [text.strip()] if text.strip() else []
    chunks = []
    for i in range(0, len(sentences), max_sentences):
        batch = sentences[i : i + max_sentences]
        chunks.append(" ".join(batch))
    logger.info("chunk_by_sentences: %s sentences -> %s chunks", len(sentences), len(chunks))
    return chunks


def chunk_course_content(course_dict: dict) -> list[dict]:
    """
    Creates intelligent chunks from a course document.
    Chunk 1: title + description (overview). Chunk 2-N: subchapters.
    """
    if not course_dict or not isinstance(course_dict, dict):
        logger.warning("chunk_course_content: invalid input")
        return []
    course_id = str(course_dict.get("id") or course_dict.get("_id", ""))
    if not course_id:
        logger.warning("chunk_course_content: no course id")
        return []
    title = course_dict.get("title") or "Untitled Course"
    description = course_dict.get("description") or ""
    sub_chapters = course_dict.get("subChapters") or course_dict.get("subchapters") or []
    if not sub_chapters:
        sub_chapters = course_dict.get("modules") or []
    if not isinstance(sub_chapters, list):
        sub_chapters = [sub_chapters]

    chunks_out = []
    # Overview chunk
    overview_text = f"{title}\n\n{description}".strip()
    if overview_text:
        chunks_out.append({
            "chunk_id": f"{course_id}_chunk_0",
            "course_id": course_id,
            "chunk_type": "overview",
            "text": overview_text,
            "metadata": {
                "course_title": title,
                "module_name": None,
                "chunk_index": 0,
            },
        })

    # One chunk per subchapter (prefer real extracted content from linked files).
    min_module_chars = 0
    buffer = []
    buffer_names = []
    chunk_index = len(chunks_out)

    for i, mod in enumerate(sub_chapters):
        if isinstance(mod, dict):
            mod_text = mod.get("title", "") or ""
            if mod.get("description"):
                mod_text += "\n" + str(mod["description"])
            real_text = _extract_subchapter_real_text(mod)
            if real_text:
                mod_text += "\n\n" + real_text
            mod_name = mod.get("title") or f"Subchapter {i + 1}"
        else:
            mod_text = str(mod)
            mod_name = f"Subchapter {i + 1}"
        buffer.append(mod_text.strip())
        buffer_names.append(mod_name)
        combined = "\n\n".join(buffer)
        if len(combined) >= min_module_chars or i == len(sub_chapters) - 1:
            if combined.strip():
                chunk_id = f"{course_id}_chunk_{chunk_index}"
                chunks_out.append({
                    "chunk_id": chunk_id,
                    "course_id": course_id,
                    "chunk_type": "module",
                    "text": combined,
                    "metadata": {
                        "course_title": title,
                        "module_name": ", ".join(buffer_names),
                        "chunk_index": chunk_index,
                    },
                })
                chunk_index += 1
            buffer = []
            buffer_names = []

    if buffer and buffer_names:
        combined = "\n\n".join(buffer)
        if combined.strip():
            chunks_out.append({
                "chunk_id": f"{course_id}_chunk_{chunk_index}",
                "course_id": course_id,
                "chunk_type": "module",
                "text": combined,
                "metadata": {
                    "course_title": title,
                    "module_name": ", ".join(buffer_names),
                    "chunk_index": chunk_index,
                },
            })

    logger.info("chunk_course_content: course_id=%s -> %s chunks", course_id, len(chunks_out))
    return chunks_out


def chunk_exercise_content(exercise_dict: dict) -> dict:
    """
    Creates a single focused chunk for an exercise (question + type + difficulty).
    """
    if not exercise_dict or not isinstance(exercise_dict, dict):
        logger.warning("chunk_exercise_content: invalid input")
        return {}
    content = exercise_dict.get("content") or exercise_dict.get("question") or ""
    ex_type = exercise_dict.get("type") or "MCQ"
    difficulty = exercise_dict.get("difficulty") or "medium"
    course_id = str(exercise_dict.get("courseId") or "")
    ex_id = str(exercise_dict.get("_id") or exercise_dict.get("id") or "")
    chunk_id = f"ex_{ex_id}" if ex_id else "ex_unknown"
    text_parts = [f"Question: {content}", f"Type: {ex_type}", f"Difficulty: {difficulty}"]
    if exercise_dict.get("options"):
        text_parts.append("Options: " + ", ".join(str(o) for o in exercise_dict["options"]))
    text = "\n".join(text_parts)
    return {
        "chunk_id": chunk_id,
        "course_id": course_id,
        "chunk_type": "exercise",
        "text": text,
        "metadata": {
            "course_title": None,
            "module_name": None,
            "chunk_index": 0,
            "exercise_type": ex_type,
            "difficulty": difficulty,
            "exercise_id": ex_id,
        },
    }


def smart_chunk_document(
    doc_type: str,
    document: dict,
    chunk_size: int = 500,
) -> list[dict]:
    """
    Dispatcher: course -> chunk_course_content, exercise -> chunk_exercise_content,
    text -> chunk_text_recursive (wrapped in chunk dicts).
    """
    try:
        if doc_type == "course":
            return chunk_course_content(document)
        if doc_type == "exercise":
            ch = chunk_exercise_content(document)
            return [ch] if ch else []
        if doc_type == "text":
            text = document.get("text", "") if isinstance(document, dict) else str(document)
            raw = chunk_text_recursive(text, chunk_size=chunk_size, chunk_overlap=50)
            return [
                {
                    "chunk_id": f"text_chunk_{i}",
                    "course_id": document.get("course_id", ""),
                    "chunk_type": "text",
                    "text": t,
                    "metadata": {"chunk_index": i},
                }
                for i, t in enumerate(raw)
            ]
        logger.warning("smart_chunk_document: unknown doc_type=%s", doc_type)
        return []
    except Exception as e:
        logger.error("smart_chunk_document error: %s", e)
        return []


def calculate_chunk_stats(chunks: list) -> dict:
    """Returns total_chunks, avg_chunk_size, min_size, max_size."""
    if not chunks:
        return {"total_chunks": 0, "avg_chunk_size": 0, "min_size": 0, "max_size": 0}
    sizes = []
    for c in chunks:
        if isinstance(c, dict):
            sizes.append(len((c.get("text") or "")))
        else:
            sizes.append(len(str(c)))
    return {
        "total_chunks": len(chunks),
        "avg_chunk_size": round(sum(sizes) / len(sizes), 1) if sizes else 0,
        "min_size": min(sizes) if sizes else 0,
        "max_size": max(sizes) if sizes else 0,
    }


def preview_chunks(chunks: list, n: int = 3) -> None:
    """Prints first N chunks for inspection."""
    for i, c in enumerate(chunks[:n]):
        text = (c.get("text") if isinstance(c, dict) else str(c))[:200]
        print(f"--- Chunk {i + 1} ---")
        print(text)
        if isinstance(c, dict) and c.get("metadata"):
            print("Metadata:", c["metadata"])
        print()


def test_chunking() -> None:
    """Runs tests with sample course and exercise documents."""
    print("=== document_chunker tests ===\n")

    # Sample course
    sample_course = {
        "id": "course_test_001",
        "title": "Introduction to Python",
        "description": "Learn Python basics. Variables, loops, and functions.",
        "subChapters": [
            {"title": "Variables", "description": "Assigning and using variables."},
            {"title": "Loops", "description": "For and while loops."},
            {"title": "Functions", "description": "Defining and calling functions."},
        ],
    }
    course_chunks = chunk_course_content(sample_course)
    print("Course chunks:", len(course_chunks))
    print("Stats:", calculate_chunk_stats(course_chunks))
    preview_chunks(course_chunks, n=2)

    # Sample exercise
    sample_ex = {
        "id": "ex_001",
        "courseId": "course_test_001",
        "content": "What is the output of print(2+2)?",
        "type": "MCQ",
        "difficulty": "easy",
        "options": ["3", "4", "5"],
    }
    ex_chunk = chunk_exercise_content(sample_ex)
    print("Exercise chunk:", ex_chunk.get("chunk_id"), len(ex_chunk.get("text", "")))

    # Recursive text
    long_text = "First paragraph.\n\nSecond paragraph here.\n\nThird one."
    rec = chunk_text_recursive(long_text, chunk_size=30, chunk_overlap=5)
    print("Recursive chunks:", len(rec), rec)

    # Sentence chunking
    sent = chunk_by_sentences("One. Two. Three. Four. Five. Six.", max_sentences=2)
    print("Sentence chunks:", sent)

    # Smart dispatcher
    smart_course = smart_chunk_document("course", sample_course)
    smart_ex = smart_chunk_document("exercise", sample_ex)
    print("Smart course chunks:", len(smart_course))
    print("Smart exercise chunks:", len(smart_ex))
    print("\nTests completed.")


if __name__ == "__main__":
    test_chunking()
